#!/usr/bin/env python3
"""Compare a modeled SWMM hydrograph with observed data."""

from __future__ import annotations

import argparse
import logging
from pathlib import Path
from typing import Iterable, Tuple

import pandas as pd

try:  # pragma: no cover - exercised at runtime only
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # type: ignore
except Exception:  # pragma: no cover - noqa: BLE001
    plt = None  # type: ignore

from .compare_outfiles import extract_series as extract_swmm_series


# ---------------------------------------------------------------------------
# Metrics
# ---------------------------------------------------------------------------

def nash_sutcliffe(sim: pd.Series, obs: pd.Series) -> float:
    """Return the Nash–Sutcliffe efficiency for ``sim`` vs ``obs``."""
    denom = ((obs - obs.mean()) ** 2).sum()
    if denom == 0:
        return float("nan")
    return 1.0 - ((obs - sim) ** 2).sum() / denom


def rmse(sim: pd.Series, obs: pd.Series) -> float:
    """Return the root mean square error for ``sim`` vs ``obs``."""
    if len(sim) == 0:
        return float("nan")
    return float(((sim - obs) ** 2).mean() ** 0.5)


def total_volume(series: pd.Series, index: pd.Index) -> float:
    """Compute the total volume under ``series`` given timestamps ``index``."""
    if len(index) < 2:
        return float("nan")
    idx = pd.to_datetime(index)
    dt = pd.Series(idx).diff().dt.total_seconds().median()
    if pd.isna(dt) or dt <= 0:
        return float("nan")
    return float(series.sum() * dt)


def add_plot_slide(pptx_path: Path, image_path: Path) -> None:
    """Append ``image_path`` as a slide in ``pptx_path``."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        logging.debug("pptx library not available; skipping PPTX export")
        return
    try:
        if pptx_path.exists():
            prs = Presentation(str(pptx_path))
        else:
            prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(0.5), width=Inches(9))
        prs.save(str(pptx_path))
        logging.info(f"Wrote PowerPoint -> {pptx_path}")
    except Exception as e:  # pragma: no cover - debug info only
        logging.error(f"Failed to write PPTX: {e}")


# ---------------------------------------------------------------------------
# CLI helpers
# ---------------------------------------------------------------------------

def parse_args(argv: Iterable[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.ArgumentDefaultsHelpFormatter
    )
    parser.add_argument("outfile", help="SWMM .out file containing modeled results")
    parser.add_argument("observed_file", help="Observed-data file (CSV or TSF)")
    parser.add_argument("element_id", help="Element ID to extract from the model")
    parser.add_argument("plot_path", help="Output path for overlay plot (PNG, etc.)")
    parser.add_argument(
        "--summary",
        default="comparison_summary.csv",
        help="Output CSV path for summary metrics",
    )
    parser.add_argument(
        "--pptx",
        default="",
        help="Path to PowerPoint file to collect plots",
    )
    parser.add_argument(
        "--item-type",
        default="node",
        choices=["node", "link", "subcatchment", "system"],
        help="SWMM object type",
    )
    parser.add_argument(
        "--param", default="Flow_rate", help="SWMM parameter name to extract"
    )
    parser.add_argument(
        "--obs-time-col",
        default="Datetime",
        help=(
            "Timestamp column in observed data; for multiple columns, provide a"
            " comma-separated list"
        ),
    )
    parser.add_argument(
        "--obs-value-col", default="Observed", help="Value column in observed data"
    )
    parser.add_argument(
        "--model-label", default="Modeled", help="Label for modeled series in plot"
    )
    parser.add_argument(
        "--obs-label", default="Observed", help="Label for observed series in plot"
    )
    parser.add_argument(
        "-v", "--verbose", action="store_true", help="Increase logging verbosity"
    )
    parser.add_argument(
        "-q", "--quiet", action="store_true", help="Suppress informational logs"
    )
    return parser.parse_args(list(argv) if argv is not None else None)


TIME_CANDIDATES = [
    "datetime",
    "date_time",
    "date/time",
    "timestamp",
    "date",
    "time",
]

VALUE_KEYWORDS = [
    "flow",
    "discharge",
    "depth",
    "velocity",
    "stage",
    "value",
]


def _guess_columns(df: pd.DataFrame, time_col: str, value_col: str) -> tuple[str, str]:
    """Return likely time and value column names for ``df``."""
    cols_l = {c.lower(): c for c in df.columns}
    # Time column
    if time_col.lower() in cols_l:
        time_col = cols_l[time_col.lower()]
    else:
        for cand in TIME_CANDIDATES:
            if cand in cols_l:
                time_col = cols_l[cand]
                break
        else:
            for col in df.columns:
                if "date" in col.lower() or "time" in col.lower():
                    time_col = col
                    break

    # Value column
    if value_col.lower() in cols_l:
        value_col = cols_l[value_col.lower()]
    else:
        for key in VALUE_KEYWORDS:
            for col in df.columns:
                if key in col.lower():
                    value_col = col
                    break
            else:
                continue
            break
    return time_col, value_col


def load_observed(
    path: str, time_col: str, value_col: str
) -> Tuple[pd.DataFrame, str, str]:
    """Load observed data from *path* which may be CSV or TSF.

    Returns a DataFrame along with the time and value column names used.
    ``time_col`` may contain a comma-separated list of columns which will be
    combined into a single ``Datetime`` column.  The column names are matched
    case-insensitively and common variants such as ``Flow_mgd`` or ``datetime``
    are recognised automatically.
    """
    ext = Path(path).suffix.lower()
    if ext == ".tsf":
        with open(path, "r", encoding="utf-8") as f:
            lines = f.readlines()
        skiprows = [2] if len(lines) > 2 and lines[2].split("\t")[0].strip() == "M/d/yyyy" else None
        df = pd.read_csv(path, sep="\t", header=1, skiprows=skiprows)

        if value_col not in df.columns and "Flow" in df.columns:
            value_col = "Flow"
    else:
        df = pd.read_csv(path)

    time_cols = [c.strip() for c in time_col.split(",") if c.strip()]
    if len(time_cols) == 1:
        tc = time_cols[0]
        if tc not in df.columns and "Date/Time" in df.columns:
            tc = "Date/Time"
        time_col = tc
    else:
        tc = [c for c in time_cols if c in df.columns]
        if len(tc) != len(time_cols):
            # missing columns; return as-is and let caller handle error
            return df, "", value_col
        subset = df[tc].copy()
        subset.columns = [c.lower() for c in subset.columns]
        try:
            df["Datetime"] = pd.to_datetime(subset)
        except Exception:
            df["Datetime"] = pd.to_datetime(subset.astype(str).agg(" ".join, axis=1))
        time_col = "Datetime"

    time_col, value_col = _guess_columns(df, time_col, value_col)

    return df, time_col, value_col


# ---------------------------------------------------------------------------
# Main workflow
# ---------------------------------------------------------------------------

def main(argv: Iterable[str] | None = None) -> int:
    args = parse_args(argv)

    log_level = logging.INFO
    if args.quiet:
        log_level = logging.WARNING
    if args.verbose:
        log_level = logging.DEBUG
    logging.basicConfig(level=log_level, format="%(levelname)s: %(message)s")

    model = extract_swmm_series(args.outfile, args.item_type, args.element_id, args.param)
    if model.empty:
        logging.error("No data extracted from model")
        return 1
    model["Datetime"] = pd.to_datetime(model["Datetime"])
    model = model.rename(columns={"value": "modeled"}).sort_values("Datetime")

    obs_df, time_col, value_col = load_observed(
        args.observed_file, args.obs_time_col, args.obs_value_col
    )
    if time_col not in obs_df.columns or value_col not in obs_df.columns:
        logging.error("Observed file missing required columns")
        return 1
    obs_df[time_col] = pd.to_datetime(obs_df[time_col])
    obs_df = obs_df.rename(columns={time_col: "Datetime", value_col: "observed"}).sort_values(
        "Datetime"
    )

    merged = pd.merge(obs_df, model, on="Datetime", how="inner")
    if merged.empty:
        logging.error("No overlapping timestamps between observed and modeled data")
        return 1

    obs = merged["observed"]
    sim = merged["modeled"]
    nse = nash_sutcliffe(sim, obs)
    rmse_val = rmse(sim, obs)
    obs_vol = total_volume(obs, merged["Datetime"])
    sim_vol = total_volume(sim, merged["Datetime"])
    if obs_vol != 0 and not pd.isna(obs_vol):
        vol_err = (sim_vol - obs_vol) / obs_vol
    else:
        vol_err = float("nan")

    if plt is None:  # pragma: no cover - handled in runtime
        logging.error("matplotlib not available; cannot create plot")
        return 1
    fig, ax = plt.subplots()
    ax.plot(merged["Datetime"], obs, label=args.obs_label)
    ax.plot(merged["Datetime"], sim, label=args.model_label)
    ax.set_xlabel("Datetime")
    ax.set_ylabel(args.param)
    ax.legend()
    fig.autofmt_xdate()
    fig.tight_layout()
    fig.savefig(args.plot_path)
    plt.close(fig)
    if args.pptx:
        add_plot_slide(Path(args.pptx), Path(args.plot_path))

    summary = pd.DataFrame(
        {
            "metric": [
                "nse",
                "rmse",
                "total_volume_error",
                "observed_volume",
                "modeled_volume",
                "count",
            ],
            "value": [nse, rmse_val, vol_err, obs_vol, sim_vol, len(merged)],
        }
    )
    summary.to_csv(args.summary, index=False)
    logging.info(f"Wrote plot to {args.plot_path} and summary to {args.summary}")
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
