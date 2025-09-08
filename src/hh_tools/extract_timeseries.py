#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SWMM Unified Extractor — v4 (Callbacks + Planning)
--------------------------------------------------
- Keeps v3 features (params discovery, units, templates, combine modes, presets)
- Adds optional progress callbacks for GUI
- Adds *planning* helpers to preview output filenames without writing
"""

from __future__ import annotations

import argparse
import glob
import json
import os
import re
import sys
import logging
from collections import defaultdict
from datetime import datetime
from typing import Dict, List, Tuple, Iterable, Optional, Any, Callable

from tqdm import tqdm

# Import swmmtoolbox and its bundled utilities eagerly so that packaging tools
# (e.g. PyInstaller) can detect the dependency tree.  Failing early if the
# vendored ``toolbox_utils`` package is missing helps avoid silent extraction
# failures.
try:  # pragma: no cover - only executed at import time
    import swmmtoolbox
    # Ensure the nested toolbox_utils tree is imported so bundlers pick it up.
    from swmmtoolbox.toolbox_utils.src.toolbox_utils import tsutils as _tsutils  # noqa: F401
except Exception as exc:  # pragma: no cover - tested indirectly
    raise ImportError(
        "swmmtoolbox with toolbox_utils is required for extract_timeseries"
    ) from exc

# Heavy deps are imported lazily where needed to keep --help fast.

# ---------------------------
# Small helpers & data types
# ---------------------------

def parse_kv_map(s: str) -> Dict[str, str]:
    out: Dict[str, str] = {}
    s = (s or "").strip()
    if not s:
        return out
    for part in s.split(","):
        part = part.strip()
        if not part:
            continue
        if "=" not in part:
            raise ValueError(f"Expected key=value, got '{part}'")
        k, v = part.split("=", 1)
        out[k.strip()] = v.strip()
    return out

def normalize_unit(u: str) -> str:
    return (u or "").strip().lower()

INVALID_FS_CHARS = set('<>:"/\\|?*')

def sanitize_id(x: str) -> str:
    """Return ``x`` made safe for filenames.

    Windows file systems forbid a small set of characters (e.g. ``*`` or
    ``?``) but allow ``#``.  We replace only the disallowed characters with
    underscores so IDs like ``Pump#1`` are preserved in the output filenames
    while still being safe to write to disk.
    """
    s = str(x)
    return "".join("_" if c in INVALID_FS_CHARS else c for c in s)

def build_output_name(pattern: str, out_format: str, **kwargs) -> str:
    """Render ``pattern`` with ``kwargs`` and ensure the extension matches ``out_format``.

    Prior to this change, any trailing ``.ext`` segment was stripped from the
    rendered pattern.  This caused IDs containing periods (e.g., ``Pump.1``) to
    lose everything after the first dot.  Now, we only remove an existing
    extension when it already matches ``out_format`` so that literal dots in IDs
    are preserved.
    """
    base = (pattern or "{prefix}{short}_{id}{suffix}").format(**kwargs)
    m = re.search(r"\.([^./]+)$", base)
    if m and m.group(1).lower() == out_format.lower():
        base = base[: m.start()]
    return f"{base}.{out_format}"

# --------------------------
# Unit conversion registers
# --------------------------

FLOW_TO_CFS = {
    "cfs": 1.0, "cms": 35.3146667, "mgd": 1.54722865, "gpm": 0.00222800926,
    "l/s": 0.0353146667,
}
LENGTH_TO_FT = {
    "ft": 1.0, "m": 3.2808399, "in": 1/12.0, "cm": 0.032808399,
}
VEL_TO_FTPS = {
    "ft/s": 1.0, "m/s": 3.2808399,
}
DIMENSIONS = {"flow","depth","head","velocity","other"}
DEFAULT_PARAM_DIM = {
    # conservative defaults (override via --param-dimension)
    "Flow_rate": "flow",
    "Flow_depth": "depth",
    "Depth": "depth",
    "Head": "head",
    "Velocity": "velocity",
    "Pump_status": "other",
}

def convert_series_by_dim(values, dim: str, from_unit: str, to_unit: str):
    import pandas as pd
    if not isinstance(values, pd.Series):
        values = pd.Series(values)

    if not from_unit or not to_unit or normalize_unit(from_unit) == normalize_unit(to_unit):
        return values

    fu = normalize_unit(from_unit)
    tu = normalize_unit(to_unit)

    if dim == "flow":
        if fu not in FLOW_TO_CFS or tu not in FLOW_TO_CFS:
            return values
        return values * (FLOW_TO_CFS[fu] / FLOW_TO_CFS[tu])

    if dim in ("depth", "head"):
        if fu not in LENGTH_TO_FT or tu not in LENGTH_TO_FT:
            return values
        v_ft = values * LENGTH_TO_FT[fu]
        return v_ft / LENGTH_TO_FT[tu]

    if dim == "velocity":
        if fu not in VEL_TO_FTPS or tu not in VEL_TO_FTPS:
            return values
        v_ftps = values * VEL_TO_FTPS[fu]
        return v_ftps / VEL_TO_FTPS[tu]

    return values

# ---------------------------------
# swmmtoolbox discovery helpers
# ---------------------------------

def list_possible_params(outfile: str, item_type: str) -> List[str]:
    """Return params available for ``item_type`` in an ``.out`` file."""
    try:
        if item_type == "pollutant":
            # swmmtoolbox does not expose pollutants via listvariables;
            # they only have one variable: concentration
            return ["Concentration"]

        vars_ = swmmtoolbox.listvariables(outfile)
    except Exception:
        return []

    params = {row[1] for row in vars_ if row and row[0] == item_type}
    return sorted(params)


def discover_ids(outfile: str, item_type: str) -> List[str]:
    """Return IDs for an ``item_type`` in an ``.out`` file."""
    try:
        if item_type == "pollutant":
            # Pollutant names live in the SwmmExtract object
            obj = swmmtoolbox.swmmtoolbox.SwmmExtract(outfile)
            return list(obj.names[3])

        catalog = swmmtoolbox.catalog(outfile, item_type)
    except Exception:
        return []

    ids = {row[1] for row in catalog if row and len(row) > 1}
    return sorted(ids)

# ----------------------------
# Export helpers (TSF / DAT)
# ----------------------------

def _write_with_headers(df, filename: str, header_lines: List[str], time_format: str, float_format: str, sep: str = "\t") -> None:
    import pandas as pd
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    with open(filename, "w", encoding="utf-8", newline="") as f:
        for h in header_lines:
            f.write(h.rstrip("\n") + "\n")
        # Write data
        if not isinstance(df.index, pd.DatetimeIndex):
            df = df.copy()
            df.index = pd.to_datetime(df.index)
        for ts, row in df.iterrows():
            f.write(ts.strftime(time_format))
            for v in row:
                if v is None:
                    f.write(sep)
                else:
                    try:
                        f.write(sep + (float_format % float(v)))
                    except Exception:
                        f.write(sep + str(v))
            f.write("\n")

def file_export_tsf(df, filename: str, header1: str, header2: str, time_format: str, float_format: str) -> None:
    _write_with_headers(df, filename, [header1, header2], time_format, float_format, sep="\t")

def file_export_dat(df, filename: str, header: str, time_format: str, float_format: str) -> None:
    import pandas as pd
    if not isinstance(df.index, pd.DatetimeIndex):
        first = df.columns[0]
        if first in df.columns and pd.api.types.is_datetime64_any_dtype(df[first]):
            df = df.set_index(first)
    header_lines = header.splitlines()
    if not any(line.startswith("Date/Time") for line in header_lines):
        header_lines.append("Date/Time\t" + "\t".join(df.columns))
    _write_with_headers(df, filename, header_lines, time_format, float_format, sep="\t")

def file_export_csv(df, filename: str, header: str, time_format: str, float_format: str) -> None:
    import pandas as pd
    if not isinstance(df.index, pd.DatetimeIndex):
        first = df.columns[0]
        if first in df.columns and pd.api.types.is_datetime64_any_dtype(df[first]):
            df = df.set_index(first)
    header_lines = header.splitlines()
    if not any(line.startswith("Date/Time") for line in header_lines):
        header_lines.append("Date/Time," + ",".join(df.columns))
    _write_with_headers(df, filename, header_lines, time_format, float_format, sep=",")

_TS_PAT = re.compile(r"^(\d{2}/\d{2}/\d{4} \d{2}:\d{2})([\t, ].*)?$")

def parse_data_lines(file_path: str, skip: int = 0) -> List[Tuple[datetime, str]]:
    out: List[Tuple[datetime, str]] = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        for i, line in enumerate(f):
            if i < skip:
                continue
            m = _TS_PAT.match(line.rstrip("\n"))
            if not m:
                continue
            ts = datetime.strptime(m.group(1), "%m/%d/%Y %H:%M")
            rest = (m.group(2) or "").lstrip("\t, ")
            out.append((ts, rest))
    return out

def read_ids_and_label_from_header(file_path: str) -> Tuple[List[str], str]:
    ids: List[str] = []
    label = ""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        for i, line in enumerate(f):
            if i == 0 and line.startswith("IDs:"):
                ids = [s.strip() for s in re.split(r"[\t,]", line)[1:]]
            elif i == 1 and line.startswith("Date/Time"):
                parts = [p.strip() for p in re.split(r"[\t,]", line)]
                if len(parts) > 1:
                    label = parts[1]
                break
    return ids, label

# ----------------------------
# Core extraction + callbacks
# ----------------------------

def extract_series(outfile: str, item_type: str, elem_id: str, param: str):
    """Return a pandas DataFrame(time,value) for a single series."""
    import pandas as pd, traceback
    # SWMM toolbox expects a label like: type,id,param  (id empty for system)
    label = f"{item_type},{elem_id},{param}"
    try:
        series = swmmtoolbox.extract(outfile, label)
    except Exception:
        with open("swmmtoolbox_error.log", "w", encoding="utf-8") as fh:
            fh.write(f"label={label}\noutfile={outfile}\n")
            fh.write(traceback.format_exc())
        raise
    if isinstance(series, pd.Series):
        return series.to_frame(name="value")
    elif isinstance(series, pd.DataFrame) and "value" in series.columns:
        return series[["value"]]
    else:
        # Normalize
        df = pd.DataFrame(series)
        if "value" not in df.columns and df.shape[1] >= 1:
            df.columns = ["value"] + list(df.columns[1:])
            df = df[["value"]]
        return df

def pretty_label(param: str, label_map: Dict[str, str], param_short: Dict[str, str]) -> Tuple[str, str]:
    """Return (column_label, short_token) for param based on maps."""
    return (label_map.get(param, param), param_short.get(param, param))

def apply_units(df, param: str, param_dimension: Dict[str, str], assume_units: Dict[str, str],
                to_units: Dict[str, str], unit_overrides: Dict[str, str]) -> Tuple[Any, Optional[str]]:
    """Convert df['value'] to desired units if mappings provided. Returns (df, output_unit or None)."""
    dim = param_dimension.get(param, DEFAULT_PARAM_DIM.get(param, "other"))
    from_u = unit_overrides.get(param, assume_units.get(dim, ""))
    to_u   = unit_overrides.get(param, to_units.get(dim, from_u))
    if from_u and to_u and from_u != to_u and dim in DIMENSIONS and dim != "other":
        df = df.copy()
        df["value"] = convert_series_by_dim(df["value"], dim, from_u, to_u)
        return df, to_u
    return df, None

def add_plot_slide(ppt: Any, df, title: str) -> None:
    try:
        if ppt is None:
            return
        import matplotlib.pyplot as plt
        from pptx import Presentation
        from pptx.util import Inches
        # create temp image
        import tempfile
        with tempfile.TemporaryDirectory() as td:
            img = os.path.join(td, "plot.png")
            plt.figure()
            ax = df.plot(legend=False)
            ax.set_title(title)
            plt.savefig(img, bbox_inches="tight")
            plt.close()
            slide_layout = ppt.slide_layouts[6]  # blank
            slide = ppt.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img, Inches(0.5), Inches(0.5), width=Inches(9))
    except Exception as e:
        logging.debug(f"Skipping PPT slide: {e}")

def plan_elements(outfile: str,
                  item_type: str,
                  element_ids: Iterable[str],
                  params: Iterable[str],
                  out_format: str,
                  combine_mode: str,
                  outdir_root: str,
                  prefix: str,
                  suffix: str,
                  dat_template: str,
                  tsf_template_sep: str,
                  tsf_template_com: str,
                  param_short: Dict[str,str]) -> List[str]:
    """Return the filenames that would be written for this selection (no IO)."""
    planned: List[str] = []
    out_dir = os.path.join(outdir_root, os.path.splitext(os.path.basename(outfile))[0])
    if item_type == "system":
        element_ids = ["SYSTEM"]
    for elem_id in element_ids:
        if combine_mode == "com":
            if out_format == "tsf":
                fname = (tsf_template_com or "{prefix}{type}_{id}{suffix}.tsf").format(
                    prefix=prefix, type=item_type, id=sanitize_id(elem_id), suffix=suffix
                )
                planned.append(os.path.join(out_dir, fname))
            else:
                # DAT/CSV: single file containing all parameters
                combined_short = "_".join(param_short.get(p, p) for p in params)
                pattern = dat_template or "{prefix}{type}_{id}{suffix}"
                fname = build_output_name(
                    pattern,
                    out_format,
                    prefix=prefix,
                    short=combined_short,
                    id=sanitize_id(elem_id),
                    suffix=suffix,
                    type=item_type,
                    param=combined_short,
                )
                planned.append(os.path.join(out_dir, fname))
        else:
            for p in params:
                if out_format == "tsf":
                    fname = (tsf_template_sep or "{prefix}{type}_{id}_{param}{suffix}.tsf").format(
                        prefix=prefix, type=item_type, id=sanitize_id(elem_id), param=p,
                        short=param_short.get(p, p), suffix=suffix
                    )
                    planned.append(os.path.join(out_dir, fname))
                else:
                    pattern = dat_template or "{prefix}{short}_{id}{suffix}"
                    short = param_short.get(p, p)
                    fname = build_output_name(pattern, out_format, prefix=prefix, short=short,
                                              id=sanitize_id(elem_id), suffix=suffix,
                                              type=item_type, param=p)
                    planned.append(os.path.join(out_dir, fname))
    return planned

def process_elements(
    outfile: str,
    item_type: str,
    element_ids: Iterable[str],
    params: Iterable[str],
    out_format: str,
    combine_mode: str,
    outdir_root: str,
    time_format: str,
    float_format: str,
    # naming
    prefix: str,
    suffix: str,
    dat_template: str,
    tsf_template_sep: str,
    tsf_template_com: str,
    param_short: Dict[str, str],
    label_map: Dict[str, str],
    # units
    param_dimension: Dict[str, str],
    assume_units: Dict[str, str],
    to_units: Dict[str, str],
    unit_overrides: Dict[str, str],
    show_progress: bool = True,
    ppt: Any | None = None,
    progress_callback: Optional[Callable[[int, int, Dict[str, Any]], None]] = None,
) -> Tuple[List[str], List[Tuple[str, str, str, str, str]]]:
    """Process a set of elements and write their time series to files.

    Returns:
        Tuple of (file paths written, failures list).  Each failure entry
        contains ``(outfile, item_type, element_id, param, error)``.
    """
    import pandas as pd

    written: List[str] = []
    failures: List[Tuple[str, str, str, str, str]] = []

    out_dir = os.path.join(outdir_root, os.path.splitext(os.path.basename(outfile))[0])
    os.makedirs(out_dir, exist_ok=True)

    # Ensure system has a single pseudo-ID
    if item_type == "system":
        element_ids = ["SYSTEM"]

    element_ids = list(element_ids)
    params = list(params)
    total = len(element_ids) * len(params)
    pbar = tqdm(total=total, desc=f"{item_type} elements", unit="series", disable=not show_progress)
    done = 0

    for elem_id in element_ids:
        frames: List[Tuple[Any, str, str]] = []  # (df, label, param)
        for p in params:
            try:
                df = extract_series(outfile, item_type, ("SYSTEM" if item_type == "system" else elem_id), p)
            except Exception as e:  # pragma: no cover - defensive
                logging.error(
                    f"Failed to extract {item_type} '{elem_id}' param '{p}': {e}"
                )
                failures.append((outfile, item_type, elem_id, p, str(e)))
                done += 1
                if progress_callback:
                    progress_callback(
                        done,
                        total,
                        {"file": outfile, "type": item_type, "id": elem_id, "param": p},
                    )
                pbar.update(1)
                continue

            df, out_u = apply_units(df, p, param_dimension, assume_units, to_units, unit_overrides)
            col_label, short = pretty_label(p, label_map, param_short)
            if out_u:
                col_label = f"{col_label} ({out_u})"
            frames.append((df, col_label, p))

            done += 1
            if progress_callback:
                progress_callback(done, total, {"file": outfile, "type": item_type, "id": elem_id, "param": p})
            pbar.update(1)

        if not frames:
            continue

        if combine_mode == "com":
            # Single file with multiple param columns
            if out_format == "tsf":
                left = frames[0][0].rename(columns={"value": frames[0][1]})
                for df, lab, _ in frames[1:]:
                    left = left.join(df.rename(columns={"value": lab}), how="outer")
                fname = (tsf_template_com or "{prefix}{type}_{id}{suffix}.tsf").format(
                    prefix=prefix, type=item_type, id=sanitize_id(elem_id), suffix=suffix
                )
                fpath = os.path.join(out_dir, fname)
                header2 = "Date/Time\t" + "\t".join([lab for _, lab, _ in frames])
                file_export_tsf(left, fpath, f"IDs:\t{elem_id}", header2, time_format, float_format)
                written.append(fpath)
            else:
                left = frames[0][0].rename(columns={"value": frames[0][1]})
                for df, lab, _ in frames[1:]:
                    left = left.join(df.rename(columns={"value": lab}), how="outer")
                combined_short = "_".join(param_short.get(p, p) for _, _, p in frames)
                pattern = dat_template or "{prefix}{type}_{id}{suffix}"
                fname = build_output_name(
                    pattern,
                    out_format,
                    prefix=prefix,
                    short=combined_short,
                    id=sanitize_id(elem_id),
                    suffix=suffix,
                    type=item_type,
                    param=combined_short,
                )
                fpath = os.path.join(out_dir, fname)
                sep = "," if out_format == "csv" else "\t"
                header = (
                    f"IDs:{sep}{elem_id}\n" +
                    "Date/Time" + sep + sep.join(param_short.get(p, p) for _, _, p in frames)
                )
                if out_format == "csv":
                    file_export_csv(left, fpath, header, time_format, float_format)
                else:
                    file_export_dat(left, fpath, header, time_format, float_format)
                written.append(fpath)
        else:
            # separate files per param
            for df, lab, p in frames:
                if out_format == "tsf":
                    fname = (tsf_template_sep or "{prefix}{type}_{id}_{param}{suffix}.tsf").format(
                        prefix=prefix, type=item_type, id=sanitize_id(elem_id), param=p,
                        short=param_short.get(p, p), suffix=suffix
                    )
                    fpath = os.path.join(out_dir, fname)
                    file_export_tsf(df.rename(columns={"value": lab}), fpath,
                                    f"IDs:\t{elem_id}", f"Date/Time\t{lab}", time_format, float_format)
                    written.append(fpath)
                else:
                    pattern = dat_template or "{prefix}{short}_{id}{suffix}"
                    fname = build_output_name(pattern, out_format, prefix=prefix,
                                               short=param_short.get(p, p),
                                               id=sanitize_id(elem_id), suffix=suffix,
                                               type=item_type, param=p)
                    fpath = os.path.join(out_dir, fname)
                    header = f"IDs:{',' if out_format == 'csv' else '\t'}{elem_id}\n" \
                             f"Date/Time{',' if out_format == 'csv' else '\t'}{param_short.get(p, p)}"
                    if out_format == "csv":
                        file_export_csv(df.rename(columns={"value": lab}), fpath,
                                        header, time_format, float_format)
                    else:
                        file_export_dat(df.rename(columns={"value": lab}), fpath,
                                        header, time_format, float_format)
                    written.append(fpath)

        for df, lab, _ in frames:
            add_plot_slide(ppt, df.rename(columns={"value": lab}), f"{item_type}:{elem_id} {lab}")

    pbar.close()
    return written, failures

def combine_across_files(
    new_files: List[Tuple[str, str]],
    out_format: str,
    output_dir: str,
    prefix: str = "",
    suffix: str = "",
    dat_template: str = "",
    tsf_template_sep: str = "",
) -> None:
    """Combine output files across elements by shared IDs, labels, and types.

    ``new_files`` should contain ``(item_type, path)`` pairs.  Only files with
    matching ``item_type`` **and** ID/label combinations are merged.  The
    resulting time series are concatenated vertically and sorted
    chronologically to mimic a continuous simulation spanning multiple ``.out``
    files.  Output naming respects user templates when provided.
    """

    import pandas as pd

    # Buckets keyed by (type, id, label)
    buckets: Dict[Tuple[str, str, str], List[str]] = defaultdict(list)

    for item_type, p in new_files:
        try:
            ids, label = read_ids_and_label_from_header(p)
            for i in ids:
                buckets[(item_type, i, label)].append(p)
        except Exception as e:
            logging.warning(f"Skipping combine for {p}: {e}")

    for (item_type, elem_id, label), paths in buckets.items():
        frames = []
        for fp in paths:
            try:
                # guess skip lines: TSF=2, DAT=1
                skip = 2 if out_format == "tsf" else 1
                rows = parse_data_lines(fp, skip=skip)
                if not rows:
                    continue
                idx = [ts for ts, _ in rows]
                vals = [float((re.split(r"[\t,]", rest)[0] if rest else "nan")) for _, rest in rows]
                frames.append(pd.DataFrame({"value": vals}, index=idx))
            except Exception as e:
                logging.warning(f"Combine read fail {fp}: {e}")
        if not frames:
            continue

        # Concatenate and sort chronologically.  Guard against an empty
        # ``frames`` list which would cause ``pd.concat`` to raise a
        # ``ValueError`` ("No objects to concatenate").  This scenario can occur
        # when the input files exist but contain no data rows.  Rather than
        # raising an exception and halting the GUI/CLI run, simply skip these
        # buckets.
        try:
            left = pd.concat(frames).sort_index()
        except ValueError as e:  # pragma: no cover - defensive
            logging.warning(
                f"Combine concat fail for {item_type} '{elem_id}' label '{label}' from {paths}: {e}"
            )
            continue
        left = left[~left.index.duplicated(keep="first")]

        out_dir = os.path.join(output_dir, "combined")
        os.makedirs(out_dir, exist_ok=True)

        short = re.sub(r"[^A-Za-z0-9]+", "_", label)
        if out_format == "tsf":
            pattern = tsf_template_sep or "{prefix}{type}_{id}_{param}{suffix}"
        else:
            pattern = dat_template or "{prefix}{short}_{id}{suffix}"
        fname = build_output_name(
            pattern,
            out_format,
            prefix=prefix,
            short=short,
            id=sanitize_id(elem_id),
            suffix=suffix,
            type=item_type,
            param=short,
        )
        out_path = os.path.join(out_dir, fname)

        if out_format == "tsf":
            file_export_tsf(
                left.rename(columns={"value": label}),
                out_path,
                f"IDs:\t{elem_id}",
                f"Date/Time\t{label}",
                "%m/%d/%Y %H:%M",
                "%.6f",
            )
        elif out_format == "csv":
            file_export_csv(
                left.rename(columns={"value": label}),
                out_path,
                f"IDs:,{elem_id}\nDate/Time,{label}",
                "%m/%d/%Y %H:%M",
                "%.6f",
            )
        else:
            file_export_dat(
                left.rename(columns={"value": label}),
                out_path,
                f"IDs:\t{elem_id}\nDate/Time\t{label}",
                "%m/%d/%Y %H:%M",
                "%.6f",
            )

def args_to_preset(args: argparse.Namespace) -> Dict[str, Any]:
    return {
        "files": args.files,
        "elements": args.elements,
        "types": args.types,
        "ids": args.ids,
        "include": args.include,
        "exclude": args.exclude,
        "want_all": args.all,
        "node_params": args.node_params,
        "link_params": args.link_params,
        "subcatchment_params": args.subcatchment_params,
        "system_params": args.system_params,
        "pollutant_params": args.pollutant_params,
        "out_format": args.out_format,
        "combine": args.combine,
        "output_dir": args.output_dir,
        "prefix": args.prefix,
        "suffix": args.suffix,
        "dat_template": args.dat_template,
        "tsf_template_sep": args.tsf_template_sep,
        "tsf_template_com": args.tsf_template_com,
        "param_short": args.param_short,
        "label_map": args.label_map,
        "assume_units": args.assume_units,
        "to_units": args.to_units,
        "unit_overrides": args.unit_overrides,
        "param_dimension": args.param_dimension,
        "time_format": args.time_format,
        "float_format": args.float_format,
        "raw": args.raw,
    }

def merge_preset(defaults: argparse.Namespace, preset: Dict[str, Any], args: argparse.Namespace) -> argparse.Namespace:
    d = vars(defaults).copy()
    d.update({k: v for k, v in preset.items() if v not in (None, "", [], {})})
    # Explicit CLI overrides final
    for k, v in vars(args).items():
        if v not in (None, "", [], {}):
            d[k] = v
    return argparse.Namespace(**d)

# ---------------------------
# CLI (unchanged behavior)
# ---------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Extract time series from SWMM .out files (v4).")
    p.add_argument("files", nargs="+", help="SWMM .out files (globs OK)")

    # Types / elements
    p.add_argument("--elements", choices=["node","link","subcatchment","system","pollutant","both"], default="both")
    p.add_argument("--types", default="", help="Comma list of types to override --elements")
    p.add_argument("--ids", default="", help="Comma list like 'node:ID1,link:ID2' (or just IDs if --elements != both)")
    p.add_argument("--include", default="", help="Regex include filter for IDs")
    p.add_argument("--exclude", default="", help="Regex exclude filter for IDs")
    p.add_argument("--all", action="store_true", help="Use all available IDs (after include/exclude)")

    # Parameter sets
    p.add_argument("--node-params", default="Flow_depth,Head")
    p.add_argument("--link-params", default="Flow_rate,Flow_depth")
    p.add_argument("--subcatchment-params", default="Runoff_flow")
    p.add_argument("--system-params", default="total_runoff_flow")
    p.add_argument("--pollutant-params", default="")

    # Output format/mode
    p.add_argument("--out-format", choices=["tsf", "dat", "csv"], default="tsf")
    p.add_argument("--combine", choices=["sep","com","across"], default="sep")
    p.add_argument("--output-dir", default="", help="Directory to write output files (defaults to input location)")
    p.add_argument("--pptx", default="", help="Path to PowerPoint file for generated plots")

    # Naming / templates
    p.add_argument("--prefix", default="", help="Prefix for output filenames")
    p.add_argument("--suffix", default="", help="Suffix for output filenames (before extension)")
    p.add_argument("--dat-template", default="", help="DAT pattern, e.g. '{prefix}{short}_{id}{suffix}.dat'")
    p.add_argument("--tsf-template-sep", default="", help="TSF per-param pattern, e.g. '{prefix}{type}_{id}_{param}{suffix}.tsf'")
    p.add_argument("--tsf-template-com", default="", help="TSF combined pattern, e.g. '{prefix}{type}_{id}{suffix}.tsf'")
    p.add_argument("--param-short", default="", help="Comma list param=SHORT for naming")
    p.add_argument("--label-map", default="", help="Comma list param=Pretty Label to override column headers")

    # Units
    p.add_argument("--assume-units", default="flow=cfs,depth=ft,head=ft,velocity=ft/s")
    p.add_argument("--to-units", default="")
    p.add_argument("--unit-overrides", default="")
    p.add_argument("--param-dimension", default="")

    # Format
    p.add_argument("--time-format", default="%m/%d/%Y %H:%M")
    p.add_argument("--float-format", default="%.6f")

    # Presets
    p.add_argument("--load-preset", default="", help="JSON preset file")
    p.add_argument("--save-preset", default="", help="Write effective preset to JSON")
    p.add_argument("--print-preset", action="store_true", help="Print the effective preset to stdout")

    # Raw
    p.add_argument("--raw", default="", help="Comma list of swmmtoolbox labels 'type,id,param' (overrides everything)")

    # Discovery
    p.add_argument("--list-params", default="", help="TYPE[,TYPE...] -> list available parameters")
    p.add_argument("--list-ids", default="", help="TYPE[,TYPE...] -> list available element IDs")
    p.add_argument("--quiet", action="store_true", help="Suppress non-error output")
    p.add_argument("-v", "--verbose", action="store_true", help="Enable debug logging")
    return p

def main(argv: Optional[List[str]] = None) -> None:
    argv = argv or sys.argv[1:]
    args = build_parser().parse_args(argv)

    level = logging.INFO
    if args.verbose:
        level = logging.DEBUG
    elif args.quiet:
        level = logging.ERROR
    logging.basicConfig(level=level, format="%(message)s")

    # Expand globs
    filelist: List[str] = []
    for patt in args.files:
        hits = sorted(glob.glob(patt))
        if not hits:
            logging.warning(f"No files matched '{patt}'")
        filelist.extend(hits)
    if not filelist:
        logging.error("No input files.")
        sys.exit(2)

    # Preset merge
    if args.load_preset:
        try:
            with open(args.load_preset, "r", encoding="utf-8") as f:
                preset = json.load(f)
            args = merge_preset(args, preset, args)
        except Exception as e:
            logging.error(f"Failed to load preset: {e}")
            sys.exit(2)

    # Save preset of effective config (before execution)
    if args.save_preset:
        try:
            with open(args.save_preset, "w", encoding="utf-8") as f:
                json.dump(args_to_preset(args), f, indent=2)
            logging.info(f"Wrote preset -> {args.save_preset}")
        except Exception as e:
            logging.error(f"Failed to save preset: {e}")

    if args.print_preset:
        print(json.dumps(args_to_preset(args), indent=2))

    if args.list_ids:
        targets = [s.strip() for s in args.list_ids.split(",") if s.strip()]
        for t in targets:
            ids = discover_ids(filelist[0], t)
            logging.info(f"[{t}] IDs:")
            for i_ in ids:
                logging.info(f"  - {i_}")
        sys.exit(0)

    # Discovery path
    if args.list_params:
        targets = [s.strip() for s in args.list_params.split(",") if s.strip()]
        for t in targets:
            params = list_possible_params(filelist[0], t)
            logging.info(f"[{t}] parameters:")
            for p_ in params:
                logging.info(f"  - {p_}")
        sys.exit(0)

    # Determine active types
    active_types = [t.strip() for t in (args.types.split(",") if args.types else []) if t.strip()]
    if not active_types:
        if args.elements == "both":
            active_types = ["node","link"]
        else:
            active_types = [args.elements]

    # Parse params per type
    params_by_type = {
        "node": [s.strip() for s in args.node_params.split(",") if s.strip()],
        "link": [s.strip() for s in args.link_params.split(",") if s.strip()],
        "subcatchment": [s.strip() for s in args.subcatchment_params.split(",") if s.strip()],
        "system": [s.strip() for s in args.system_params.split(",") if s.strip()],
        "pollutant": [s.strip() for s in args.pollutant_params.split(",") if s.strip()],
    }

    # Parse maps
    param_short = parse_kv_map(args.param_short)
    label_map   = parse_kv_map(args.label_map)
    assume_units = parse_kv_map(args.assume_units)
    to_units     = parse_kv_map(args.to_units)
    unit_overrides = parse_kv_map(args.unit_overrides)
    param_dimension = parse_kv_map(args.param_dimension)

    # Handle raw → bypass selection
    if args.raw.strip():
        new_files: List[Tuple[str, str]] = []
        all_failures: List[Tuple[str, str, str, str, str]] = []
        ppt = None
        if args.pptx:
            try:
                from pptx import Presentation
                ppt = Presentation()
            except Exception as e:
                logging.warning(f"PPTX disabled: {e}")

        labels = [s.strip() for s in args.raw.split(",") if s.strip()]
        total = len(filelist) * len(labels)
        pbar = tqdm(total=total, unit="series", disable=args.quiet, desc="extract")

        def cb(done, tot, ctx):
            pbar.update(1)

        for outfile in filelist:
            outdir_root = args.output_dir or os.path.dirname(outfile)
            for label in labels:
                try:
                    itype, elem_id, param = [part.strip() for part in label.split(",")]
                except Exception:
                    logging.error(f"Bad raw label: {label}")
                    continue
                written, failures = process_elements(
                    outfile=outfile,
                    item_type=itype,
                    element_ids=[elem_id],
                    params=[param],
                    out_format=args.out_format,
                    combine_mode="sep",
                    outdir_root=outdir_root,
                    time_format=args.time_format,
                    float_format=args.float_format,
                    prefix=args.prefix,
                    suffix=args.suffix,
                    dat_template=args.dat_template,
                    tsf_template_sep=args.tsf_template_sep,
                    tsf_template_com=args.tsf_template_com,
                    param_short=param_short,
                    label_map=label_map,
                    param_dimension=param_dimension,
                    assume_units=assume_units,
                    to_units=to_units,
                    unit_overrides=unit_overrides,
                    show_progress=False,
                    ppt=ppt,
                    progress_callback=cb,
                )
                new_files.extend((itype, f) for f in written)
                all_failures.extend(failures)

        pbar.close()

        if args.combine == "across" and new_files:
            combine_across_files(
                new_files,
                args.out_format,
                (args.output_dir or os.getcwd()),
                prefix=args.prefix,
                suffix=args.suffix,
                dat_template=args.dat_template,
                tsf_template_sep=args.tsf_template_sep,
            )
        if ppt and args.pptx:
            try:
                ppt.save(args.pptx)
                logging.info(f"Wrote PowerPoint -> {args.pptx}")
            except Exception as e:
                logging.error(f"Failed to save PowerPoint: {e}")
        if all_failures:
            lines = ["The following elements could not be exported:"]
            for f, t, i, p, err in all_failures:
                lines.append(f"- {os.path.basename(f)} [{t}] {i} ({p}): {err}")
            logging.warning("\n".join(lines))
        logging.info("Done.")
        return

    # Expanded selection path
    new_files: List[Tuple[str, str]] = []
    all_failures: List[Tuple[str, str, str, str, str]] = []
    ppt = None
    if args.pptx:
        try:
            from pptx import Presentation
            ppt = Presentation()
        except Exception as e:
            logging.warning(f"PPTX disabled: {e}")

    # Pre-compute IDs per file to establish total progress
    per_file_ids: List[Tuple[str, Dict[str, List[str]]]] = []
    total = 0
    for outfile in filelist:
        ids_by_type: Dict[str, List[str]] = {}
        if args.all or not args.ids.strip():
            for t in active_types:
                present = discover_ids(outfile, t)
                flt = []
                inc = re.compile(args.include) if args.include else None
                exc = re.compile(args.exclude) if args.exclude else None
                for i in present:
                    if inc and not inc.search(i):
                        continue
                    if exc and exc.search(i):
                        continue
                    flt.append(i)
                ids_by_type[t] = flt
        else:
            ids_by_type = defaultdict(list)
            for token in [s.strip() for s in args.ids.split(",") if s.strip()]:
                if ":" in token:
                    t, idv = token.split(":", 1)
                else:
                    if len(active_types) != 1:
                        logging.error(f"Ambiguous ID '{token}' — specify as type:ID")
                        sys.exit(2)
                    t, idv = active_types[0], token
                ids_by_type[t].append(idv)
        per_file_ids.append((outfile, ids_by_type))

        for item_type in active_types:
            element_ids = ids_by_type.get(item_type, [])
            params = params_by_type.get(item_type, [])
            if item_type == "system" and element_ids:
                element_ids = ["SYSTEM"]
            total += (len(element_ids) or 0) * (len(params) or 0)

    total = max(total, 1)
    pbar = tqdm(total=total, unit="series", disable=args.quiet, desc="extract")

    def cb(done, tot, ctx):
        pbar.update(1)

    for outfile, ids_by_type in per_file_ids:
        outdir_root = args.output_dir or os.path.dirname(outfile)
        for item_type in active_types:
            element_ids = ids_by_type.get(item_type, [])
            if not element_ids:
                continue
            params = params_by_type.get(item_type, [])
            if not params:
                continue
            written, failures = process_elements(
                outfile=outfile,
                item_type=item_type,
                element_ids=element_ids,
                params=params,
                out_format=args.out_format,
                combine_mode=args.combine,
                outdir_root=outdir_root,
                time_format=args.time_format,
                float_format=args.float_format,
                prefix=args.prefix,
                suffix=args.suffix,
                dat_template=args.dat_template,
                tsf_template_sep=args.tsf_template_sep,
                tsf_template_com=args.tsf_template_com,
                param_short=param_short,
                label_map=label_map,
                param_dimension=param_dimension,
                assume_units=assume_units,
                to_units=to_units,
                unit_overrides=unit_overrides,
                show_progress=False,
                ppt=ppt,
                progress_callback=cb,
            )
            new_files.extend((item_type, f) for f in written)
            all_failures.extend(failures)

    pbar.close()

    if args.combine == "across" and new_files:
        combine_across_files(
            new_files,
            args.out_format,
            (args.output_dir or os.getcwd()),
            prefix=args.prefix,
            suffix=args.suffix,
            dat_template=args.dat_template,
            tsf_template_sep=args.tsf_template_sep,
        )

    if ppt and args.pptx:
        try:
            ppt.save(args.pptx)
            logging.info(f"Wrote PowerPoint -> {args.pptx}")
        except Exception as e:
            logging.error(f"Failed to save PowerPoint: {e}")

    if all_failures:
        lines = ["The following elements could not be exported:"]
        for f, t, i, p, err in all_failures:
            lines.append(f"- {os.path.basename(f)} [{t}] {i} ({p}): {err}")
        logging.warning("\n".join(lines))

    logging.info("Done.")

if __name__ == "__main__":
    main()
